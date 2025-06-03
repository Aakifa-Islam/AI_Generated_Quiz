import io 
import docx
import fitz  # PyMuPDF
from pptx import Presentation


def extract_text_from_pdf_content(content_bytes):
    text = ""
    # Open the PDF from the content bytes using io.BytesIO
    doc = fitz.open(stream=content_bytes, filetype="pdf")
    
    # Loop through the pages and extract text
    for page in doc:
        text += page.get_text()

    return text




# Function to extract text from DOCX content bytes
def extract_text_from_docx_content(content_bytes):
    doc = docx.Document(io.BytesIO(content_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

# Function to extract text from PPTX content bytes
def extract_text_from_pptx_content(content_bytes):
    text = ""
    prs = Presentation(io.BytesIO(content_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to extract text from TXT content bytes
def extract_text_from_txt_content(content_bytes):
    return content_bytes.decode("utf-8")







# Function to check file type and extract content
def extract_text_from_content(content_bytes, file_extension):
    ext = file_extension.lower()
    if ext == "pdf":
        return extract_text_from_pdf_content(content_bytes)
    elif ext == "docx":
        return extract_text_from_docx_content(content_bytes)
    elif ext == "pptx":
        return extract_text_from_pptx_content(content_bytes)
    elif ext == "txt":
        return extract_text_from_txt_content(content_bytes)
    else:
        return "Unsupported file type"










